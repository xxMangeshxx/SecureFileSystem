import os
import hashlib
from cryptography.fernet import Fernet
from PyQt5 import QtWidgets, QtCore
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QPushButton, QLabel, QFileDialog, QMessageBox, QVBoxLayout, QDialog, QProgressBar, QInputDialog, QPlainTextEdit
)
import base64
import zipfile
import io
import base64
from PyQt5.QtCore import QThread, pyqtSignal
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

UPLOADS_FOLDER = "uploads"
ENCRYPTED_FOLDER = "encrypted"
KEY_FILE = "key.key"
ACTIONS_FILE = "actions.txt"

class LoginPage(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Admin Login")
        self.setGeometry(200, 200, 400, 200)
        self.user_role = None  # Initialize user role here
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout()

        self.username_label = QLabel("Username:", self)
        self.username_input = QtWidgets.QLineEdit(self)

        self.password_label = QLabel("Password:", self)
        self.password_input = QtWidgets.QLineEdit(self)
        self.password_input.setEchoMode(QtWidgets.QLineEdit.Password)

        self.login_button = QPushButton("Login", self)
        self.login_button.clicked.connect(self.login)

        layout.addWidget(self.username_label)
        layout.addWidget(self.username_input)
        layout.addWidget(self.password_label)
        layout.addWidget(self.password_input)
        layout.addWidget(self.login_button)

        self.setLayout(layout)

    def login(self):
        username = self.username_input.text()
        password = self.password_input.text()

        if username == "admin" and password == "admin123":
            self.user_role = "admin"
            self.accept()  # Proceed to the main page as admin
        elif username == "user" and password == "user123":
            self.user_role = "user"
            self.accept()  # Proceed to the main page as user
        else:
            QMessageBox.warning(self, "Login Failed", "Invalid credentials. Please try again.")

    def get_user_role(self):
        return self.user_role  # Return the user role after login

class FileSharingServer:
    def __init__(self):
        if not os.path.exists(UPLOADS_FOLDER):
            os.makedirs(UPLOADS_FOLDER)
        if not os.path.exists(ENCRYPTED_FOLDER):
            os.makedirs(ENCRYPTED_FOLDER)
        if not os.path.exists(KEY_FILE):
            key = Fernet.generate_key()
            with open(KEY_FILE, "wb") as key_file:
                key_file.write(key)
        with open(KEY_FILE, "rb") as key_file:
            self.key = key_file.read()
            self.cipher = Fernet(self.key)

        if not os.path.exists(ACTIONS_FILE):
            with open(ACTIONS_FILE, "w") as file:
                file.write("Previous actions:\n")
        else:
            with open(ACTIONS_FILE, "w") as file:
                file.write("Previous actions:\n")

    def encrypt_file(self, file_name, data):
        encrypted_data = self.cipher.encrypt(data)
        encrypted_file_path = os.path.join(ENCRYPTED_FOLDER, file_name)
        with open(encrypted_file_path, "wb") as file:
            file.write(encrypted_data)
        return encrypted_file_path

    def decrypt_file(self, file_name):
        encrypted_file_path = os.path.join(ENCRYPTED_FOLDER, file_name)
        with open(encrypted_file_path, "rb") as file:
            encrypted_data = file.read()

        # Decrypt the data
        decrypted_data = self.cipher.decrypt(encrypted_data)

        return decrypted_data


    def upload_file(self, file_name, data, compress=False):
    # Encrypt the file data first
        encrypted_data = self.encrypt_file(file_name, data)

        # If compression is enabled, compress the encrypted data
        if compress:
            encrypted_data = self.compress_file(encrypted_data)

        file_path = os.path.join(UPLOADS_FOLDER, file_name)
        if os.path.exists(file_path):
            return None
        with open(file_path, "wb") as file:
            file.write(encrypted_data)

        # Log the action of uploading
        with open(ACTIONS_FILE, "a") as actions_file:
            actions_file.write(f"Uploaded file: {file_name}\n")
        return file_name


    def compress_file(self, data):
        try:
            compressed_data = io.BytesIO()
            with zipfile.ZipFile(compressed_data, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                zip_file.writestr("data", data)  # You can name the file inside the zip appropriately
            compressed_data.seek(0)
            return compressed_data.read()
        except Exception as e:
            print(f"Error compressing file: {e}")
            return None

    def decompress_file(self, data):
        try:
            # Check if the data is a valid zip file (starts with b'PK')
            if not data.startswith(b'PK'):
                print("Data is not a valid zip file.")
                return None

            with zipfile.ZipFile(io.BytesIO(data), 'r') as zip_ref:
                file_name = zip_ref.namelist()[0]  # Get the file name from the zip
                decompressed_data = zip_ref.read(file_name)  # Read the file content inside the zip
            return decompressed_data
        except zipfile.BadZipFile as e:
            print(f"Bad zip file: {e}")
            return None
        except Exception as e:
            print(f"Error decompressing file: {str(e)}")
            return None
    def list_files(self):
        return os.listdir(UPLOADS_FOLDER)

    def list_encrypted_files(self):
        return os.listdir(ENCRYPTED_FOLDER)

    def get_previous_actions(self):
        with open(ACTIONS_FILE, "r") as actions_file:
            return actions_file.read()


class MainWindow(QMainWindow):
    def __init__(self, server, user_role="user"):
        super().__init__()
        self.server = server
        self.user_role = user_role
        self.setWindowTitle("File Sharing Application")
        self.setGeometry(200, 200, 800, 600)
        self.init_ui()

    def init_ui(self):
        self.title_label = QLabel("File Sharing Application", self)
        self.title_label.setStyleSheet("font-size: 24px; font-weight: bold;")
        self.title_label.setAlignment(QtCore.Qt.AlignCenter)

        self.upload_button = QPushButton("Upload File", self)
        self.upload_button.setStyleSheet("font-size: 18px; padding: 10px;")
        self.upload_button.clicked.connect(self.upload_file)

        self.download_button = QPushButton("Download File", self)
        self.download_button.setStyleSheet("font-size: 18px; padding: 10px;")
        self.download_button.clicked.connect(self.download_file)

        self.list_button = QPushButton("List Files", self)
        self.list_button.setStyleSheet("font-size: 18px; padding: 10px;")
        self.list_button.clicked.connect(self.list_files)

        self.view_actions_button = QPushButton("View Previous Actions", self)
        self.view_actions_button.setStyleSheet("font-size: 18px; padding: 10px;")
        self.view_actions_button.clicked.connect(self.view_previous_actions)

        if self.user_role == "admin":
            self.view_encrypted_button = QPushButton("View Encrypted Files", self)
            self.view_encrypted_button.setStyleSheet("font-size: 18px; padding: 10px;")
            self.view_encrypted_button.clicked.connect(self.view_encrypted_files)
            self.view_encrypted_button.setVisible(True)

        layout = QVBoxLayout()
        layout.addWidget(self.title_label)
        layout.addWidget(self.upload_button)
        layout.addWidget(self.download_button)
        layout.addWidget(self.list_button)
        layout.addWidget(self.view_actions_button)

        if self.user_role == "admin":
            layout.addWidget(self.view_encrypted_button)

        central_widget = QtWidgets.QWidget()
        central_widget.setLayout(layout)
        self.setCentralWidget(central_widget)

    def upload_file(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Select File to Upload")
        if file_path:
            try:
                compress = True  # Set to True to compress file before upload
                file_extension = os.path.splitext(file_path)[1].lower()

                # Show original file size
                original_size = os.path.getsize(file_path)
                print(f"Original file size: {original_size / 1024:.2f} KB")

                # Read file contents based on file extension
                if file_extension == '.pdf':
                    file_contents = self.read_pdf(file_path)
                elif file_extension == '.docx':
                    file_contents = self.read_word(file_path)
                elif file_extension == '.pptx':
                    file_contents = self.read_ppt(file_path)
                elif file_extension in ['.jpg', '.jpeg', '.png']:
                    file_contents = self.read_image(file_path)
                else:
                    file_contents = self.read_text(file_path)

                # Display file contents in a message box or text edit
                if file_contents:
                    self.show_file_contents(file_contents)

                with open(file_path, "rb") as file:
                    file_data = file.read()

                # Conditionally compress file if it's large enough
                if compress and original_size > 1024:  # Only compress if the file is larger than 1 KB
                    compressed_data = self.server.compress_file(file_data)
                    compressed_size = len(compressed_data)
                    print(f"Compressed file size: {compressed_size / 1024:.2f} KB")
                else:
                    compressed_data = file_data  # No compression for small files

                # Prepare file name and upload
                file_name = os.path.basename(file_path)
                upload_successful = self.server.upload_file(file_name, compressed_data, compress)

                if upload_successful:
                    QMessageBox.information(self, "Success", f"File '{file_name}' uploaded successfully!")
                else:
                    QMessageBox.warning(self, "Failed", f"File '{file_name}' already exists.")
            except Exception as e:
                QMessageBox.warning(self, "Error", f"An error occurred: {str(e)}")

    def show_file_contents(self, contents):
        # Create a QDialog to show the file contents
        dialog = QDialog(self)
        dialog.setWindowTitle("File Contents")
        dialog.setGeometry(300, 300, 600, 400)
        
        # Create a QPlainTextEdit widget to display contents
        text_edit = QPlainTextEdit(dialog)
        text_edit.setPlainText(contents)
        text_edit.setReadOnly(True)  # Make it read-only
        text_edit.setGeometry(10, 10, 580, 380)

        # Add a button to close the dialog
        close_button = QPushButton("Close", dialog)
        close_button.setGeometry(250, 350, 100, 30)
        close_button.clicked.connect(dialog.accept)

        dialog.exec_()

    def download_file(self):
        files = self.server.list_files()
        if files:
            file_name, ok = QtWidgets.QInputDialog.getItem(self, "Select File", "Files on server:", files, 0, False)
            if ok and file_name:
                destination_folder = QFileDialog.getExistingDirectory(self, "Select Destination Folder")
                if destination_folder:
                    try:
                        encrypted_file_path = os.path.join(UPLOADS_FOLDER, file_name)

                        # Read encrypted file
                        with open(encrypted_file_path, "rb") as file:
                            encrypted_data = file.read()

                        # Decrypt the data
                        decrypted_data = self.server.decrypt_file(file_name)

                        # Decompress if it was compressed during upload
                        final_data = self.server.decompress_file(decrypted_data) or decrypted_data

                        # Save the decrypted and decompressed file
                        file_path = os.path.join(destination_folder, file_name)
                        with open(file_path, "wb") as file:
                            file.write(final_data)

                        QMessageBox.information(self, "Success", f"File '{file_name}' downloaded and saved at: {file_path}")
                    except Exception as e:
                        QMessageBox.warning(self, "Error", f"An error occurred during download: {str(e)}")





    def list_files(self):
        files = self.server.list_files()
        if files:
            files_info = "\n".join(files)
        else:
            files_info = "No files available"

        # Debug: Check what's in the UPLOADS_FOLDER
        print(f"Files in upload folder: {files_info}")

        QMessageBox.information(self, "Files List", files_info)


    def view_previous_actions(self):
        actions = self.server.get_previous_actions()
        QMessageBox.information(self, "Previous Actions", actions)

    def view_encrypted_files(self):
        encrypted_files = self.server.list_encrypted_files()
        if encrypted_files:
            file_name, ok = QtWidgets.QInputDialog.getItem(self, "Select Encrypted File", "Encrypted files:", encrypted_files, 0, False)
            if ok and file_name:
                try:
                    decrypted_data = self.server.decrypt_file(file_name)

                    # Optionally, decompress the decrypted data if it's compressed
                    decompressed_data = self.server.decompress_file(decrypted_data)

                    # Display the decrypted data as base64 encoded string
                    if decompressed_data:
                        encoded_content = base64.b64encode(decompressed_data).decode()  # Encode the decompressed data
                        self.show_file_contents(encoded_content)  # Display the encoded content in the dialog
                    else:
                        QMessageBox.warning(self, "Decompression Failed", "Failed to decompress the file.")
                except Exception as e:
                    QMessageBox.warning(self, "Error", f"An error occurred: {str(e)}")
        else:
            QMessageBox.information(self, "No Encrypted Files", "No encrypted files available.")


    def read_pdf(self, file_path):
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text

    def read_word(self, file_path):
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text
        return text

    def read_ppt(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        return text

    def read_image(self, file_path):
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text

    def read_text(self, file_path):
        with open(file_path, 'r') as file:
            return file.read()

class LoadingScreen(QDialog):
    def __init__(self, text):
        super().__init__()
        self.setWindowTitle("Loading...")
        self.setGeometry(200, 200, 300, 100)
        layout = QVBoxLayout()
        self.progress_bar = QProgressBar(self)
        self.progress_bar.setRange(0, 0)  # Indeterminate progress bar
        self.label = QLabel(text, self)
        layout.addWidget(self.label)
        layout.addWidget(self.progress_bar)
        self.setLayout(layout)

def main():
    app = QApplication([])
    login = LoginPage()

    if login.exec_() == QDialog.Accepted:
        user_role = login.get_user_role()
        server = FileSharingServer()
        window = MainWindow(server, user_role)
        window.show()
        app.exec_()

if __name__ == "__main__":
    main()
